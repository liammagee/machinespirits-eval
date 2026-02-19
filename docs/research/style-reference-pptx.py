#!/usr/bin/env python3
"""Style the pandoc reference.pptx to match the Geist in the Machine theme."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
import copy

# ─── Color Palette (matches slides-header.tex) ──────────────
GEIST_PRIMARY = RGBColor(0x1B, 0x28, 0x38)   # Dark slate
GEIST_ACCENT  = RGBColor(0xD4, 0x87, 0x2C)   # Warm amber
GEIST_LIGHT   = RGBColor(0xF5, 0xF2, 0xEB)   # Warm off-white
GEIST_MID     = RGBColor(0x5C, 0x6B, 0x7A)   # Medium slate
GEIST_TEXT    = RGBColor(0x2D, 0x34, 0x36)    # Near-black
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEADING = 'Helvetica Neue'
FONT_BODY    = 'Helvetica Neue'
FONT_MONO    = 'Fira Mono'


def set_slide_bg(slide_layout, color):
    """Set background color on a slide layout."""
    bg = slide_layout.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_placeholder(ph, font_name=None, font_size=None, font_color=None,
                       bold=None, alignment=None):
    """Style a placeholder's default text properties."""
    if ph.has_text_frame:
        for paragraph in ph.text_frame.paragraphs:
            if alignment is not None:
                paragraph.alignment = alignment
            for run in paragraph.runs:
                if font_name:
                    run.font.name = font_name
                if font_size:
                    run.font.size = font_size
                if font_color:
                    run.font.color.rgb = font_color
                if bold is not None:
                    run.font.bold = bold

    # Also set the default text style via XML for new text
    sp = ph._sp
    txBody = sp.find(qn('p:txBody'))
    if txBody is not None:
        for defRPr_parent in txBody.findall(qn('a:lstStyle')):
            for level in defRPr_parent:
                defRPr = level.find(qn('a:defRPr'))
                if defRPr is not None and font_name:
                    # Set latin font
                    latin = defRPr.find(qn('a:latin'))
                    if latin is not None:
                        latin.set('typeface', font_name)
                    else:
                        latin = defRPr.makeelement(qn('a:latin'), {'typeface': font_name})
                        defRPr.append(latin)


def update_theme_colors(prs):
    """Update the theme color scheme to match our palette."""
    from lxml import etree

    # Access the slide master's theme
    slide_master = prs.slide_masters[0]

    # Find theme XML through the relationship
    for rel in slide_master.part.rels.values():
        if 'theme' in rel.reltype:
            theme_element = etree.fromstring(rel.target_part.blob)
            theme_part_ref = rel.target_part

            # Find the color scheme
            theme_elements = theme_element.findall('.//' + qn('a:clrScheme'))
            for clr_scheme in theme_elements:
                # Update specific theme colors
                color_map = {
                    'dk1': GEIST_PRIMARY,    # Dark 1
                    'dk2': GEIST_TEXT,       # Dark 2
                    'lt1': WHITE,            # Light 1
                    'lt2': GEIST_LIGHT,      # Light 2
                    'accent1': GEIST_ACCENT, # Accent 1
                    'accent2': GEIST_MID,    # Accent 2
                    'accent3': RGBColor(0x27, 0xAE, 0x60),  # Green
                    'accent4': RGBColor(0xC0, 0x39, 0x2B),  # Red
                    'hlink': GEIST_ACCENT,   # Hyperlink
                    'folHlink': GEIST_MID,   # Followed hyperlink
                }

                for color_name, rgb in color_map.items():
                    el = clr_scheme.find(qn(f'a:{color_name}'))
                    if el is not None:
                        # Remove existing color children
                        for child in list(el):
                            el.remove(child)
                        # Add srgbClr
                        srgb = el.makeelement(qn('a:srgbClr'), {'val': str(rgb)})
                        el.append(srgb)

            # Update font scheme
            font_schemes = theme_element.findall('.//' + qn('a:fontScheme'))
            for font_scheme in font_schemes:
                for font_type in ['majorFont', 'minorFont']:
                    font_el = font_scheme.find(qn(f'a:{font_type}'))
                    if font_el is not None:
                        latin = font_el.find(qn('a:latin'))
                        if latin is not None:
                            typeface = FONT_HEADING if font_type == 'majorFont' else FONT_BODY
                            latin.set('typeface', typeface)

            # Save modified theme back
            theme_part_ref._blob = etree.tostring(theme_element, xml_declaration=True,
                                                   encoding='UTF-8', standalone=True)
            break


def main():
    prs = Presentation('reference.pptx')

    # Update theme colors
    update_theme_colors(prs)

    # Style each slide layout
    slide_master = prs.slide_masters[0]

    for layout in slide_master.slide_layouts:
        layout_name = layout.name.lower()

        if 'title' in layout_name and 'content' not in layout_name:
            # Title slide - dark background
            set_slide_bg(layout, GEIST_PRIMARY)

            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:  # Title
                    style_placeholder(ph, font_name=FONT_HEADING,
                                     font_size=Pt(36), font_color=WHITE,
                                     bold=True)
                elif ph.placeholder_format.idx == 1:  # Subtitle
                    style_placeholder(ph, font_name=FONT_BODY,
                                     font_size=Pt(18), font_color=GEIST_ACCENT)

        elif 'section' in layout_name:
            # Section header - dark background
            set_slide_bg(layout, GEIST_PRIMARY)

            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:  # Title
                    style_placeholder(ph, font_name=FONT_HEADING,
                                     font_size=Pt(32), font_color=WHITE,
                                     bold=True)
                else:
                    style_placeholder(ph, font_name=FONT_BODY,
                                     font_color=GEIST_ACCENT)

        elif 'two' in layout_name and 'content' in layout_name:
            # Two-column layout
            set_slide_bg(layout, WHITE)
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:  # Title
                    style_placeholder(ph, font_name=FONT_HEADING,
                                     font_size=Pt(28), font_color=GEIST_PRIMARY,
                                     bold=True)
                else:
                    style_placeholder(ph, font_name=FONT_BODY,
                                     font_size=Pt(16), font_color=GEIST_TEXT)

        elif 'blank' in layout_name:
            set_slide_bg(layout, WHITE)

        else:
            # Content slides - white background
            set_slide_bg(layout, WHITE)

            for ph in layout.placeholders:
                idx = ph.placeholder_format.idx
                if idx == 0:  # Title
                    style_placeholder(ph, font_name=FONT_HEADING,
                                     font_size=Pt(28), font_color=GEIST_PRIMARY,
                                     bold=True)
                    # Add bottom border to title area via shape properties
                elif idx == 1:  # Body/content
                    style_placeholder(ph, font_name=FONT_BODY,
                                     font_size=Pt(18), font_color=GEIST_TEXT)

    # Style the slide master itself
    set_slide_bg(slide_master, WHITE)

    prs.save('reference.pptx')
    print('Styled reference.pptx successfully.')


if __name__ == '__main__':
    main()
